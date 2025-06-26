import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
import pptx
import pandas as pd
import os
import time
from langchain.chains.retrieval_qa.base import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_openai import ChatOpenAI
from langchain_google_genai import ChatGoogleGenerativeAI
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import re

st.set_page_config(
    page_title="SDLC Automate APP",
    page_icon="images/favicon.png"
)

# Get the API key from Streamlit secrets
os.environ["OPENAI_API_KEY"] = st.secrets["OPENAI_API_KEY"]
os.environ["GOOGLE_API_KEY"] = st.secrets["GOOGLE_API_KEY"]

# Function to calculate confidence level based on prompt-answer accuracy
def calculate_confidence_level(prompt, answer):
    """Calculate confidence level based on how well the answer addresses the prompt."""
    try:
        # Preprocess texts
        prompt_clean = re.sub(r'[^\w\s]', '', prompt.lower())
        answer_clean = re.sub(r'[^\w\s]', '', answer.lower())
        
        # Extract key terms from prompt
        prompt_keywords = set(prompt_clean.split())
        answer_words = set(answer_clean.split())
        
        # Calculate keyword overlap
        common_keywords = prompt_keywords.intersection(answer_words)
        keyword_overlap = len(common_keywords) / len(prompt_keywords) if prompt_keywords else 0
        
        # Use TF-IDF for semantic similarity
        vectorizer = TfidfVectorizer(stop_words='english', max_features=1000)
        tfidf_matrix = vectorizer.fit_transform([prompt_clean, answer_clean])
        similarity_score = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0]
        
        # Combined confidence score (weighted average)
        confidence_score = (keyword_overlap * 0.3 + similarity_score * 0.7) * 100
        
        return min(confidence_score, 100)  # Cap at 100%
    except Exception as e:
        st.error(f"Error calculating confidence: {str(e)}")
        return 0

# Function to calculate match percentage with source document
def calculate_match_percentage(answer, source_text):
    """Calculate how well the answer matches the source document content."""
    try:
        # Preprocess texts
        answer_clean = re.sub(r'[^\w\s]', '', answer.lower())
        source_clean = re.sub(r'[^\w\s]', '', source_text.lower())
        
        # Use TF-IDF for document similarity
        vectorizer = TfidfVectorizer(stop_words='english', max_features=1000, ngram_range=(1, 2))
        tfidf_matrix = vectorizer.fit_transform([source_clean, answer_clean])
        match_score = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0]
        
        # Convert to percentage
        match_percentage = match_score * 100
        
        return min(match_percentage, 100)  # Cap at 100%
    except Exception as e:
        st.error(f"Error calculating match percentage: {str(e)}")
        return 0

# Function to get confidence level category
def get_confidence_category(percentage):
    """Convert percentage to confidence level category."""
    if percentage >= 70:
        return "High", "🟢"
    elif percentage >= 30:
        return "Medium", "🟡"
    else:
        return "Low", "🔴"

# Function to display confidence metrics
def display_confidence_metrics(confidence_score, match_score):
    """Display confidence and match metrics in a formatted way."""
    conf_level, conf_icon = get_confidence_category(confidence_score)
    match_level, match_icon = get_confidence_category(match_score)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.metric(
            label="Confidence Level",
            value=f"{conf_level} ({confidence_score:.1f}%)",
            help="How well the answer addresses the prompt"
        )
        st.write(f"{conf_icon} **Confidence Score:** {confidence_score:.1f}%")
    
    with col2:
        st.metric(
            label="Match Percentage", 
            value=f"{match_level} ({match_score:.1f}%)",
            help="How well the answer matches the source document"
        )
        st.write(f"{match_icon} **Match Score:** {match_score:.1f}%")
    
    # Overall assessment
    overall_score = (confidence_score + match_score) / 2
    overall_level, overall_icon = get_confidence_category(overall_score)
    
    st.info(f"{overall_icon} **Overall Assessment:** {overall_level} ({overall_score:.1f}%)")

# Streamlit sidebar setup
with st.sidebar:
    st.title("Your BRD Documents")
    model_selection = st.selectbox(
        "Select AI Model",
        options=["Open AI GPT 4o", "Google Gemini 2.0 Flash"]
    )
    st.write(f"Selected Model: {model_selection}")
    uploaded_file = st.file_uploader("Upload a file to generate user stories", type=["pdf", "docx", "txt", "xlsx", "pptx"])
    
    # Add confidence settings
    st.subheader("Quality Assessment Settings")
    show_confidence = st.checkbox("Show Confidence & Match Analysis", value=True)
    detailed_metrics = st.checkbox("Show Detailed Metrics", value=False)

# Function to extract text from various file types
@st.cache_resource
def extract_text_from_file(file):
    """Extracts text based on file type, with caching for faster retrieval."""
    text = ""
    file_ext = os.path.splitext(file.name)[1].lower()

    # Handle PDF files
    if file_ext == ".pdf":
        pdf_reader = PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()

    # Handle Word (.docx) files
    elif file_ext == ".docx":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"

    # Handle text (.txt) files
    elif file_ext == ".txt":
        text = file.read().decode("utf-8")

    # Handle Excel files (.xlsx, .xls)
    elif file_ext in [".xlsx", ".xls"]:
        df = pd.read_excel(file)
        text = df.to_string()

    # Handle PowerPoint files (.pptx, .ppt)
    elif file_ext in [".pptx", ".ppt"]:
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    return text

# Process the uploaded file and extract text for the vector store
@st.cache_resource
def process_uploaded_file(uploaded_file):
    return extract_text_from_file(uploaded_file) if uploaded_file else ""

# Function to create vector store from extracted text
@st.cache_resource
def create_vector_store(text):
    text_splitter = RecursiveCharacterTextSplitter(
        separators="\n",
        chunk_size=800,
        chunk_overlap=50,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    embeddings = OpenAIEmbeddings()
    return FAISS.from_texts(chunks, embeddings)

# Streamlit app setup
st.header("BRD to User Story, Test Case, Cucumber Script, and Selenium Script")

# Set up tabs for different functionalities
tab1, tab2, tab3, tab4 = st.tabs(["User Story Generation", "User Story to Test Case", "Test Case to Cucumber Script", "Test Case to Selenium Script"])

# Initialize variables
text = ""
vector_store = None
qa_chain = None

# Process uploaded file if available
if uploaded_file:
    text = process_uploaded_file(uploaded_file)
    if text:
        vector_store = create_vector_store(text)
        
        if model_selection == "Open AI GPT 4o":
            llm = ChatOpenAI(
                model="gpt-4o",
                temperature=0.7,
            )
        elif model_selection == "Google Gemini 2.0 Flash":
            llm = ChatGoogleGenerativeAI(
                model="gemini-2.0-flash",
                temperature=0.7,
            )
        
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=vector_store.as_retriever()
        )

# User Story Generation Tab
with tab1:
    start_time = time.time()
    if uploaded_file and qa_chain:
        prompt_message = (
            """ Imagine you are a Senior Business Analyst. "

                "Your responsibility is to read the entire Business Requirement Document (BRD) "

                "and convert it into detailed User Stories. "

                "Think step-by-step and ensure you write all possible User Stories derived from the BRD." 

                "Provide fully complete User Stories only, "

                "without any additional explanation or sentences."

                "give only user stories with the standard format of writing user stories """
        )
        
        start_query_time = time.time()
        matches = vector_store.similarity_search(prompt_message, k=3)
        response = qa_chain.invoke({"query": prompt_message})
        
        st.subheader("Generated User Stories")
        st.write(response['result'])
        
        # Calculate and display confidence metrics
        if show_confidence:
            st.subheader("Quality Assessment")
            
            confidence_score = calculate_confidence_level(prompt_message, response['result'])
            match_score = calculate_match_percentage(response['result'], text)
            
            display_confidence_metrics(confidence_score, match_score)
            
            if detailed_metrics:
                with st.expander("Detailed Analysis"):
                    st.write("**Prompt Analysis:**")
                    st.write(f"- Prompt length: {len(prompt_message)} characters")
                    st.write(f"- Response length: {len(response['result'])} characters")
                    st.write(f"- Source document length: {len(text)} characters")
                    
                    st.write("**Retrieved Context:**")
                    for i, match in enumerate(matches):
                        st.write(f"Match {i+1}: {match.page_content[:200]}...")

        # Display timing info
        st.write(f"Document loading time: {time.time() - start_time:.2f} seconds")
        st.write(f"Query processing time: {time.time() - start_query_time:.2f} seconds")
    else:
        st.write("Please upload a BRD document in the sidebar to generate user stories.")

# User Story to Test Case Tab
with tab2:
    st.subheader("Convert User Story to Test Case")
    user_story_text = st.text_area("Enter the user story text here to generate test cases:")

    if st.button("Generate Test Cases"):
        if user_story_text and qa_chain:
            test_case_prompt = (
                "You are a senior QA engineer. Your responsibility is to design a comprehensive test suite for the following user story: \n\n" + user_story_text + 
                "\n\n**Objectives:**" + 
                "\n* **Develop** a set of test cases that cover all aspects of the user story functionality." + 
                "\n* **Consider** both functional and edge case scenarios." + 
                "\n* **Ensure** test cases are presented in a clear and concise manner." + 
                "\n* **Focus** solely on the test cases themselves, omitting any additional explanations or context." 
            )
            start_test_case_time = time.time()
            response = qa_chain.invoke({"query": test_case_prompt})
            
            st.subheader("Generated Test Cases")
            st.write(response['result'])
            
            # Calculate and display confidence metrics
            if show_confidence:
                st.subheader("Quality Assessment")
                
                confidence_score = calculate_confidence_level(test_case_prompt, response['result'])
                match_score = calculate_match_percentage(response['result'], user_story_text)
                
                display_confidence_metrics(confidence_score, match_score)
            
            st.write(f"Test case generation time: {time.time() - start_test_case_time:.2f} seconds")
        elif not qa_chain:
            st.write("Please upload a BRD document first to initialize the AI model.")
        else:
            st.write("Please enter a user story to generate test cases.")

# Test Case to Cucumber Script Tab
with tab3:
    st.subheader("Convert Test Case to Cucumber Script")
    test_case_text = st.text_area("Enter the test case text here to generate Cucumber script:")

    if st.button("Generate Cucumber Script"):
        if test_case_text and qa_chain:
            cucumber_prompt = (
                "You are a test automation engineer tasked with creating a Cucumber test suite for the following test case: \n\n" + test_case_text + 
                "\n\n**Objectives:**" + 
                "\n* **Develop** a comprehensive Cucumber feature file using Gherkin syntax." + 
                "\n* **Create** corresponding Java step definitions for each Gherkin step." + 
                "\n* **Ensure** all scenarios are defined with Given, When, and Then steps as appropriate." + 
                "\n* **Focus** on producing clean, readable, and maintainable code." + 
                "\n* **Deliver** both the complete feature file and the complete Java step definition file."
            )
            start_cucumber_time = time.time()
            response = qa_chain.invoke({"query": cucumber_prompt})
            
            st.subheader("Generated Cucumber Script")
            st.write(response['result'])
            
            # Calculate and display confidence metrics
            if show_confidence:
                st.subheader("Quality Assessment")
                
                confidence_score = calculate_confidence_level(cucumber_prompt, response['result'])
                match_score = calculate_match_percentage(response['result'], test_case_text)
                
                display_confidence_metrics(confidence_score, match_score)
            
            st.write(f"Cucumber script generation time: {time.time() - start_cucumber_time:.2f} seconds")
        elif not qa_chain:
            st.write("Please upload a BRD document first to initialize the AI model.")
        else:
            st.write("Please enter a test case to generate a Cucumber script.")

# Test Case to Selenium Script Tab
with tab4:
    st.subheader("Convert Test Case to Selenium Script")
    selenium_test_case_text = st.text_area("Enter the test case text here to generate Selenium script:")

    if st.button("Generate Selenium Script"):
        if selenium_test_case_text and qa_chain:
            selenium_prompt = (
                "Assume you are a test automation engineer specializing in Selenium. Your task is to convert the following test case "
                "into a Selenium WebDriver script using Python. Ensure to include all steps to perform the actions in the test case, "
                "Make sure to give fully complete selenium full code."
                "such as locating elements, interacting with the web page, and validating outcomes. Here is the test case: \n\n" + selenium_test_case_text
            )
            start_selenium_time = time.time()
            response = qa_chain.invoke({"query": selenium_prompt})
            
            st.subheader("Generated Selenium Script")
            st.write(response['result'])
            
            # Calculate and display confidence metrics
            if show_confidence:
                st.subheader("Quality Assessment")
                
                confidence_score = calculate_confidence_level(selenium_prompt, response['result'])
                match_score = calculate_match_percentage(response['result'], selenium_test_case_text)
                
                display_confidence_metrics(confidence_score, match_score)
            
            st.write(f"Selenium script generation time: {time.time() - start_selenium_time:.2f} seconds")
        elif not qa_chain:
            st.write("Please upload a BRD document first to initialize the AI model.")
        else:
            st.write("Please enter a test case to generate a Selenium script.")
